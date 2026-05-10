from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from rag.ingest.asset_anchors import asset_anchor
from rag.ingest.parsers.util import default_title_from_location, normalize_whitespace, slugify
from rag.schema.core import ParsedDocument, ParsedElement, ParsedSection, SourceType


class PptxParserRepo:
    """
    专门针对 PPTX 格式的解析特种兵。
    核心使命：原生解析、保留幻灯片二维结构、提取最关键的“演讲者备注(Speaker Notes)”。
    """
    
    def parse(
        self,
        file_path: Path,
        *,
        location: str,
        source_type: SourceType,
        title: str | None = None,
        owner: str = "user",
    ) -> ParsedDocument:
        prs = Presentation(file_path)

        # 提取原生核心元数据
        core_props = prs.core_properties
        doc_title = title or core_props.title or default_title_from_location(location)
        doc_author = core_props.author or owner

        sections: list[ParsedSection] = []
        all_elements: list[ParsedElement] = []
        visible_parts: list[str] = []
        visible_cursor = 0
        visible_section_separator = "\n\n"

        for slide_idx, slide in enumerate(prs.slides):
            page_no = slide_idx + 1
            slide_title = self._extract_slide_title(slide) or f"Slide {page_no}"

            slide_elements: list[ParsedElement] = []
            slide_text_parts: list[str] = []

            # 1. 抽取幻灯片标题
            if slide.shapes.title and slide.shapes.title.text:
                title_text = normalize_whitespace(slide.shapes.title.text)
                if title_text:
                    slide_elements.append(
                        ParsedElement(
                            element_id=f"slide-{page_no}-title",
                            kind="section_header",
                            text=title_text,
                            toc_path=(doc_title, slide_title),
                            heading_level=2,
                            page_no=page_no,
                            metadata={},
                        )
                    )
                    slide_text_parts.append(title_text)

            # 2. 抽取版面元素
            for shape_idx, shape in enumerate(slide.shapes):
                if shape == slide.shapes.title:
                    continue

                if shape.has_text_frame:
                    text = normalize_whitespace(shape.text)
                    if text:
                        slide_elements.append(
                            ParsedElement(
                                element_id=f"slide-{page_no}-shape-{shape_idx}",
                                kind="text",
                                text=text,
                                toc_path=(doc_title, slide_title),
                                page_no=page_no,
                                metadata={},
                            )
                        )
                        slide_text_parts.append(text)

                elif shape.has_table:
                    md_table = self._table_to_markdown(shape.table)
                    if md_table:
                        element_id = f"slide-{page_no}-table-{shape_idx}"
                        slide_elements.append(
                            ParsedElement(
                                element_id=element_id,
                                kind="table",
                                text=md_table,
                                toc_path=(doc_title, slide_title),
                                page_no=page_no,
                                metadata={},
                            )
                        )
                        slide_text_parts.append(asset_anchor(element_id))

            # 3. 抽取演讲者备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = normalize_whitespace(slide.notes_slide.notes_text_frame.text)
                if notes_text:
                    slide_elements.append(
                        ParsedElement(
                            element_id=f"slide-{page_no}-notes",
                            kind="speaker_note",
                            text=notes_text,
                            toc_path=(doc_title, slide_title),
                            page_no=page_no,
                            metadata={"is_speaker_note": "true"},
                        )
                    )
                    slide_text_parts.append(f"Speaker Notes: {notes_text}")

            # 4. 封装成 Section：在写入 visible_text 的同时记录 span
            section_text = normalize_whitespace(" \n".join(slide_text_parts))
            if section_text:
                if visible_parts:
                    visible_parts.append(visible_section_separator)
                    visible_cursor += len(visible_section_separator)

                char_range_start = visible_cursor
                visible_parts.append(section_text)
                visible_cursor += len(section_text)
                char_range_end = visible_cursor

                sections.append(
                    ParsedSection(
                        toc_path=(doc_title, slide_title),
                        heading_level=2,
                        page_range=(page_no, page_no),
                        order_index=slide_idx,
                        text=section_text,
                        char_range_start=char_range_start,
                        char_range_end=char_range_end,
                        anchor_hint=slugify(f"slide-{page_no}-{slide_title}"),
                        metadata={"slide_number": str(page_no)},
                    )
                )
                all_elements.extend(slide_elements)

        visible_text = "".join(visible_parts)

        for idx, section in enumerate(sections):
            start = section.char_range_start
            end = section.char_range_end
            if start is None or end is None:
                raise ValueError(f"pptx section[{idx}] missing char range")
            if start < 0 or end <= start or end > len(visible_text):
                raise ValueError(
                    f"pptx section[{idx}] invalid char range: start={start}, end={end}, visible_len={len(visible_text)}"
                )
            if visible_text[start:end] != section.text:
                raise ValueError(
                    f"pptx section[{idx}] text/span mismatch: "
                    f"expected={section.text!r}, actual={visible_text[start:end]!r}"
                )

        return ParsedDocument(
            title=doc_title,
            source_type=SourceType.PPTX,
            authors=[doc_author],
            language=None,
            sections=sections,
            visible_text=visible_text,
            visual_semantics=None,
            elements=all_elements,
            page_count=len(prs.slides),
            doc_model=None,
            metadata={"location": location, "source_type": SourceType.PPTX.value},
        )

    def _extract_slide_title(self, slide: Any) -> str | None:
        if slide.shapes.title and slide.shapes.title.text:
            return normalize_whitespace(slide.shapes.title.text)
        return None

    def _table_to_markdown(self, table: Any) -> str:
        """原生解析 PPT 内置表格，转为稳定的 Markdown"""
        if not table.rows:
            return ""
            
        md_lines = []
        for i, row in enumerate(table.rows):
            # 将单元格内的换行符替换为空格，防止打乱 Markdown 表格结构
            row_data = [normalize_whitespace(cell.text_frame.text).replace("\n", " ") for cell in row.cells]
            md_lines.append("| " + " | ".join(row_data) + " |")
            
            # 添加表头分隔符
            if i == 0:
                md_lines.append("|" + "|".join(["---"] * len(row.cells)) + "|")
                
        return "\n".join(md_lines)
